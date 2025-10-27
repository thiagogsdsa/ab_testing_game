# =============================
# Standard library imports
# =============================
import os
import sys
import warnings
warnings.filterwarnings("ignore")

# =============================
# Third-party imports
# =============================
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')  # headless backend
import matplotlib.pyplot as plt
import seaborn as sns
from scipy import stats
import subprocess
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from ab_bayes_test import ABBayesTest  

# =============================
# Local application imports
# =============================
current_dir = os.path.dirname(os.path.abspath(__file__))
project_root = os.path.abspath(os.path.join(current_dir, '..'))
sys.path.append(project_root)

from src.helper_functions import (
    plot_histograms_data_1vsdata_2,
    plot_histogram_single,
    plot_kde_data1vsdata_2_same_plot,
    simulate_profit_posterior
)
# -----------------------------
# PowerPoint Slide Utilities
# -----------------------------
def create_graph_slide(prs, title, subtitle, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tf = textbox.text_frame
    tf.clear()
    
    p_title = tf.paragraphs[0]
    run = p_title.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(32)
    run.font.color.rgb = RGBColor(0, 0, 0)
    p_title.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        p_sub = tf.add_paragraph()
        run_sub = p_sub.add_run()
        run_sub.text = subtitle
        run_sub.font.size = Pt(20)
        run_sub.font.color.rgb = RGBColor(50, 50, 50)
        p_sub.alignment = PP_ALIGN.CENTER

    # Image
    if image_path:
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.8), width=Inches(9))
    
    return slide

def build_presentation(slide_creators, output_file="presentation.pptx"):
    output_dir = os.path.dirname(output_file)
    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation()
    for func, args in slide_creators:
        func(prs, **args)
    prs.save(output_file)
    print(f"Presentation saved at: {output_file}")

# -----------------------------
# LaTeX Slide Utilities
# -----------------------------
def create_recommendation_slide_image(value_col, control_group, treatment_group, lift_summary, slides_dir="slides"):
    os.makedirs(slides_dir, exist_ok=True)
    
    prob = lift_summary['prob_treatment_superior'] * 100
    val_col_safe = value_col.replace("_", r"\_")
    treat_group_safe = treatment_group.replace("_", r"\_")
    control_group_safe = control_group.replace("_", r"\_")
    
    path_tex = os.path.join(slides_dir, f"{value_col}_slide_4.tex")
    path_pdf = os.path.join(slides_dir, f"{value_col}_slide_4.pdf")
    path_png = os.path.join(slides_dir, f"{value_col}_slide_4.png")

    latex_content = f"""\\documentclass[16pt]{{article}}
\\usepackage{{amsmath,amssymb,xcolor}}
\\usepackage[active,tightpage]{{preview}}
\\PreviewEnvironment{{center}}
\\setlength\\parindent{{0pt}}
\\pagestyle{{empty}}
\\begin{{document}}
\\begin{{center}}
\\textbf{{}}\\\\[1em]
No statistically significant improvement detected. Maintain current version.\\\\[1em]
Probability that

$$
\\text{{Prob}}(\\text{{Treatment}} > \\text{{Control}}) = {prob:.2f}\\%
$$

is far below the 95\\% threshold.
\\end{{center}}
\\end{{document}}"""
    
    with open(path_tex, "w") as f:
        f.write(latex_content)
    
    subprocess.run(["pdflatex", "-output-directory", slides_dir, path_tex], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    
    pages = convert_from_path(path_pdf, dpi=200)
    pages[0].save(path_png, "PNG")
    
    return path_png

def create_conclusion_slide_image(value_col="final", slides_dir="slides"):
    os.makedirs(slides_dir, exist_ok=True)
    
    path_tex = os.path.join(slides_dir, f"{value_col}_slide_1.tex")
    path_pdf = os.path.join(slides_dir, f"{value_col}_slide_1.pdf")
    path_png = os.path.join(slides_dir, f"{value_col}_slide_1.png")

    latex_content = rf"""\documentclass[16pt]{{article}}
\usepackage{{amsmath,amssymb,xcolor}}
\usepackage[active,tightpage]{{preview}}
\PreviewEnvironment{{center}}
\setlength\parindent{{0pt}}
\pagestyle{{empty}}
\begin{{document}}
\begin{{center}}
\textbf{{Final Recommendation}} \\[1em]
After reviewing all metrics, including retention on Day 1 and Day 7,
engagement through total rounds, and overall business performance, 
the results show no consistent evidence that the treatment version delivers better outcomes. \\[1em]
The recommendation is to keep the current control version and continue monitoring performance 
over the next cycles.
\end{{center}}
\end{{document}}"""

    with open(path_tex, "w") as f:
        f.write(latex_content)
    
    subprocess.run(["pdflatex", "-output-directory", slides_dir, path_tex], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    
    pages = convert_from_path(path_pdf, dpi=200)
    pages[0].save(path_png, "PNG")
    
    return path_png

# -----------------------------
# Data Preparation
# -----------------------------
df_raw = pd.read_csv(os.path.join(project_root, "data", "cookie_cats.csv"))
df_raw['userid'] = df_raw['userid'].astype(str)
df_checks = df_raw.copy()
df_checks['sum_gamerounds_log'] = np.log1p(df_checks['sum_gamerounds'])
df_checks = df_checks.query("sum_gamerounds > 0")

# -----------------------------
# Slide Generation
# -----------------------------
def generate_all_slide_images(
    df=df_checks,
    group_col='version',
    value_col='retention_1',
    metric_type='proportion',
    inference_type='conjugate',
    prior_params=None,
    sampling_size=10000,
    control_group='gate_30',
    treatment_group='gate_40',
    alpha: float = 0.05,
    slides_dir="slides"
):
    paths = []

    # Fit Bayesian model
    ab_test = ABBayesTest(
        df=df,
        group_col=group_col,
        value_col=value_col,
        metric_type=metric_type,
        inference_type=inference_type,
        prior_params=prior_params,
        sampling_size=sampling_size,
        control_group=control_group,
        treatment_group=treatment_group
    )
    ab_test.fit()

    # Posterior distributions and lift
    posterior_samples = ab_test.get_distributions()
    lift_samples = posterior_samples[treatment_group] / posterior_samples[control_group] - 1

    posterior_summary = ab_test.results(1 - alpha)
    lift_summary = ab_test.lift_summary()

    df_post = pd.DataFrame(posterior_summary).T.reset_index().rename(columns={'index': 'metric'})
    df_lift = pd.DataFrame([lift_summary]).T.reset_index().rename(columns={'index':'metric', 0:'value'})

    # Posterior histogram
    path1 = f"{slides_dir}/{value_col}_slide_0.png"
    plot_histograms_data_1vsdata_2(posterior_samples, path_to_save=path1)
    paths.append(path1)

    # Posterior KDE
    path2 = f"{slides_dir}/{value_col}_slide_1.png"
    plot_kde_data1vsdata_2_same_plot(posterior_samples, path_to_save=path2)
    paths.append(path2)

    # Lift distribution
    path3 = f"{slides_dir}/{value_col}_slide_2.png"
    plot_histogram_single(lift_samples, path_to_save=path3)
    paths.append(path3)

    # Summary tables
    path4 = f"{slides_dir}/{value_col}_slide_3.png"
    fig, axes = plt.subplots(2, 1, figsize=(6, 4))

    axes[0].axis('off')
    table_post = axes[0].table(
        cellText=df_post.round(4).values,
        colLabels=df_post.columns,
        cellLoc='center',
        loc='center'
    )
    table_post.auto_set_font_size(False)
    table_post.set_fontsize(9)
    table_post.scale(1.2, 1.2)

    axes[1].axis('off')
    table_lift = axes[1].table(
        cellText=df_lift.round(4).values,
        colLabels=df_lift.columns,
        cellLoc='center',
        loc='center'
    )
    table_lift.auto_set_font_size(False)
    table_lift.set_fontsize(9)
    table_lift.scale(1.2, 1.2)

    plt.tight_layout()
    plt.savefig(path4, bbox_inches='tight')
    plt.close()
    paths.append(path4)

    # Recommendation slide
    path5 = create_recommendation_slide_image(
        value_col=value_col,
        control_group=control_group,
        treatment_group=treatment_group,
        lift_summary=lift_summary
    )
    paths.append(path5)

    # Optional 7-day profit projection
    path6 = f"{slides_dir}/final_slide_0.png"
    if value_col == 'retention_7':
        simulate_profit_posterior(
            posterior_samples=posterior_samples,
            posterior_summary=posterior_summary,
            initial_players=len(df_checks),
            ltv_per_player=3,
            title="Projected 7-Day Profit Distribution",
            path_to_save=path6
        )
        paths.append(path6)

    return paths

# -----------------------------
# Generate Slides & PowerPoint
# -----------------------------
slides_to_create = {}
variables = ["retention_1", "retention_7", "sum_gamerounds_log"]

for var in variables:
    if "retention" in var:
        day = "1" if var == "retention_1" else "7"
        title_prefix = f"Retention Day {day}"
        metric_type = "proportion"
    else:
        title_prefix = "Total Rounds"
        metric_type = "mean"

    slides = generate_all_slide_images(
        df=df_checks,
        value_col=var,
        metric_type=metric_type,
        slides_dir='slides'
    )

    slides_to_create[var] = [
        (create_graph_slide, {"title": title_prefix, "subtitle":"Posterior Distribution", "image_path": slides[0]}),
        (create_graph_slide, {"title": title_prefix, "subtitle":"Posterior Distribution", "image_path": slides[1]}),
        (create_graph_slide, {"title": title_prefix, "subtitle":"Lift Distribution", "image_path": slides[2]}),
        (create_graph_slide, {"title": title_prefix, "subtitle":"Summary", "image_path": slides[3]}),
        (create_graph_slide, {"title": title_prefix, "subtitle":"Recommendations", "image_path": slides[4]})
    ]

# Conclusion slide
create_conclusion_slide_image()
slides_to_create["final"] = [
    (create_graph_slide, {"title":"Business Performance", "subtitle":"Posterior Distribution", "image_path":"slides/final_slide_0.png"}),
    (create_graph_slide, {"title":"Conclusion", "subtitle":"", "image_path":"slides/final_slide_1.png"})
]

# Combine all slides
slides_to_create_final = (
    slides_to_create["retention_1"] +
    slides_to_create["retention_7"] +
    slides_to_create["sum_gamerounds_log"] +
    slides_to_create["final"]
)

# Generate PowerPoint
build_presentation(slides_to_create_final, output_file="powerpoint/ab_test_full_presentation.pptx")

# -----------------------------
# LaTeX Beamer Report
# -----------------------------
import re

def latex_escape(s: str) -> str:
    replacements = {
        '\\': r'\textbackslash{}',
        '{': r'\{',
        '}': r'\}',
        '$': r'\$',
        '&': r'\&',
        '#': r'\#',
        '_': r'\_',
        '%': r'\%',
        '~': r'\textasciitilde{}',
        '^': r'\textasciicircum{}',
    }
    pattern = re.compile('|'.join(re.escape(key) for key in replacements.keys()))
    return pattern.sub(lambda x: replacements[x.group()], s)

def build_beamer_report_from_existing_images(slides_dict, slides_dir="pdf", output_file="ab_test_report.pdf"):
    tex_path = os.path.join(slides_dir, "ab_test_report.tex")
    os.makedirs(slides_dir, exist_ok=True)

    tex_lines = [
    r"\documentclass{beamer}",
    r"\usetheme{Singapore}",
    r"\usecolortheme{seahorse}",
    r"\useinnertheme{rectangles}",
    r"\useoutertheme{miniframes}",
    r"\usepackage{graphicx,amsmath,amssymb,xcolor,url}",
    r"\title{AB Test Report}",
    r"\author{Thiago Guimarães Santos \\ thiago.guimaraes.sto@gmail.com \\ linkedin.com/in/thiagogsdsa}",
    r"\date{\today}",
    r"\begin{document}",
    r"\frame{\titlepage}"
 ]

    for var, image_paths in slides_dict.items():
        slide_titles = [
            "Posterior Distribution",
            "Posterior KDE",
            "Lift Distribution",
            "Summary",
            "Recommendation"
        ]
        for title, img_path in zip(slide_titles, image_paths):
            tex_lines.append(rf"\begin{{frame}}{{{latex_escape(var)} - {latex_escape(title)}}}")
            if img_path.endswith(".tex"):
                tex_lines.append(rf"\input{{{img_path}}}")
            else:
                tex_lines.append(rf"\includegraphics[width=\textwidth]{{{latex_escape(img_path)}}}")
            tex_lines.append(r"\end{frame}")

    tex_lines.append(r"\end{document}")

    with open(tex_path, "w") as f:
        f.write("\n".join(tex_lines))

    subprocess.run(["pdflatex", "-output-directory", slides_dir, tex_path], check=True)

    final_pdf_path = os.path.join(slides_dir, output_file)
    print(f"PDF Beamer report generated at: {final_pdf_path}")
    return final_pdf_path


slides_dict = {
    "retention_1": [
        "slides/retention_1_slide_0.png",
        "slides/retention_1_slide_1.png",
        "slides/retention_1_slide_2.png",
        "slides/retention_1_slide_3.png",
        "slides/retention_1_slide_4.png",
        "slides/final_slide_0.png",
        "slides/final_slide_1.png"

    ],
    "retention_7": [
        "slides/retention_7_slide_0.png",
        "slides/retention_7_slide_1.png",
        "slides/retention_7_slide_2.png",
        "slides/retention_7_slide_3.png",
        "slides/retention_7_slide_4.png",
        "slides/final_slide_0.png",
        "slides/final_slide_1.png"
    ],
    "sum_gamerounds_log": [
        "slides/sum_gamerounds_log_slide_0.png",
        "slides/sum_gamerounds_log_slide_1.png",
        "slides/sum_gamerounds_log_slide_2.png",
        "slides/sum_gamerounds_log_slide_3.png",
        "slides/sum_gamerounds_log_slide_4.png",
        "slides/final_slide_0.png",
        "slides/final_slide_1.png"
    ]
}

pdf_beamer = build_beamer_report_from_existing_images(slides_dict, slides_dir="pdf_beamer")